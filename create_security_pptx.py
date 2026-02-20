"""
Generate Security Protocols PowerPoint with OpenExchange branding.
Brand colors: Teal #008285, Dark #000000, White #FFFFFF
Font: Lora (fallback to Calibri in PPTX)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
TEAL = RGBColor(0x00, 0x82, 0x85)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
BLUE = RGBColor(0x25, 0x63, 0xEB)
TEAL_LIGHT = RGBColor(0xE0, 0xF7, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=DARK, bullet_color=TEAL):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_icon_card(slide, left, top, width, height, icon_text, title, description, accent_color=TEAL):
    """Add a card with icon, title, and description."""
    # Card background
    card = add_shape(slide, left, top, width, height, WHITE, border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(1))

    # Accent bar on left
    add_shape(slide, left, top, Inches(0.06), height, accent_color)

    # Icon circle
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3), Inches(0.5), Inches(0.5))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = accent_color
    icon_shape.line.fill.background()
    tf = icon_shape.text_frame
    tf.paragraphs[0].text = icon_text
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = False

    # Title
    add_text_box(slide, left + Inches(1.0), top + Inches(0.2), width - Inches(1.2), Inches(0.4),
                 title, font_size=16, color=DARK, bold=True)

    # Description
    add_text_box(slide, left + Inches(1.0), top + Inches(0.6), width - Inches(1.2), height - Inches(0.8),
                 description, font_size=11, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)

# Teal accent bar at top
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)

# Company name
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(6), Inches(0.6),
             "OpenExchange", font_size=20, color=TEAL, bold=True)

# Title
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.2),
             "Security Protocols & Controls", font_size=44, color=WHITE, bold=True)

# Subtitle
add_text_box(slide, Inches(0.8), Inches(3.4), Inches(10), Inches(0.8),
             "Kaltura \u2192 AWS \u2192 Zoom Video Migration Platform", font_size=22, color=RGBColor(0x94, 0xA3, 0xB8))

# Divider line
add_shape(slide, Inches(0.8), Inches(4.5), Inches(2), Inches(0.04), TEAL)

# Meta info
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(0.4),
             "Confidential  |  February 2026", font_size=14, color=MID_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.4),
             "Prepared for: OpenExchange Engineering & Security Teams", font_size=14, color=MID_GRAY)

# Decorative element on right
add_shape(slide, Inches(10.5), Inches(1.5), Inches(2.5), Inches(5.5), RGBColor(0x00, 0x6B, 0x6E))
add_shape(slide, Inches(11.0), Inches(2.0), Inches(2.333), Inches(5.5), RGBColor(0x00, 0x55, 0x58))


# ═══════════════════════════════════════════════════════════
# SLIDE 2: Executive Summary
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.5),
             "Executive Summary", font_size=32, color=DARK, bold=True)

add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), TEAL)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
             "Security isn't one thing \u2014 it's identity, data protection, blast radius control, and auditability. "
             "This document outlines the security protocols required to move the migration platform from a functional demo to an enterprise-grade, sellable product.",
             font_size=16, color=MID_GRAY)

# 7 pillars
pillars = [
    ("\U0001F512", "Identity & API Security", "OAuth 2.0, IAM, Secrets Mgmt"),
    ("\U0001F6E1", "Data Protection", "Encryption at rest & in transit"),
    ("\U0001F310", "Network Security", "VPC, private endpoints, WAF"),
    ("\U0001F4CB", "Access Control", "RBAC, SSO, MFA on dashboard"),
    ("\u26A1", "Rate Limiting", "Backoff, circuit breakers"),
    ("\U0001F50D", "Audit & Logging", "CloudWatch, structured logs"),
    ("\U0001F3DB", "Compliance", "SOC 2, GDPR, HIPAA readiness"),
]

for i, (icon, title, desc) in enumerate(pillars):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + col * Inches(3.1)
    y = Inches(2.8) + row * Inches(2.0)

    card = add_shape(slide, x, y, Inches(2.8), Inches(1.6), LIGHT_GRAY, border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(1))

    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.4), Inches(0.35),
                 f"{icon}  {title}", font_size=15, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.4), Inches(0.8),
                 desc, font_size=12, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: Identity & API Security
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "1. Identity & API Security", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "The biggest risk area \u2014 securing authentication across three platforms.", font_size=14, color=MID_GRAY)

# Kaltura card
add_icon_card(slide, Inches(0.8), Inches(2.0), Inches(3.8), Inches(4.8),
              "K", "Kaltura Authentication",
              "Protocol: Server-to-server session auth\n\n"
              "\u2022 KS token via session.start API\n"
              "\u2022 24-hour TTL, auto-refresh\n"
              "\u2022 Read-only media permissions\n"
              "\u2022 No admin delete access\n"
              "\u2022 Rotate Partner Secret quarterly\n"
              "\u2022 IP restriction if supported\n\n"
              "Store Partner Secret in:\n"
              "  AWS Secrets Manager\n"
              "  NEVER in source code",
              accent_color=RGBColor(0x7C, 0x3A, 0xED))

# Zoom card
add_icon_card(slide, Inches(4.9), Inches(2.0), Inches(3.8), Inches(4.8),
              "Z", "Zoom OAuth 2.0 S2S",
              "Protocol: Server-to-Server (Client Credentials)\n\n"
              "\u2022 POST zoom.us/oauth/token\n"
              "\u2022 grant_type=account_credentials\n"
              "\u2022 Short-lived tokens (\u22641 hour)\n"
              "\u2022 No token persistence in DB\n"
              "\u2022 Scoped: clips:write only\n"
              "\u2022 Rotate every 90 days\n\n"
              "DO NOT use:\n"
              "  \u2718 User OAuth\n"
              "  \u2718 JWT apps (deprecated)\n"
              "  \u2718 Hardcoded tokens",
              accent_color=BLUE)

# AWS IAM card
add_icon_card(slide, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.8),
              "A", "AWS IAM Protocol",
              "Protocol: IAM Roles (no static keys)\n\n"
              "\u2022 Lambda role: least privilege\n"
              "\u2022 s3:PutObject on specific bucket\n"
              "\u2022 dynamodb:UpdateItem on table\n"
              "\u2022 secretsmanager:GetSecretValue\n"
              "\u2022 Deny everything else\n\n"
              "NEVER grant:\n"
              "  \u2718 AdministratorAccess\n"
              "  \u2718 S3 full access\n"
              "  \u2718 Root account usage\n\n"
              "Enable MFA for human IAM users",
              accent_color=RGBColor(0xD9, 0x77, 0x06))


# ═══════════════════════════════════════════════════════════
# SLIDE 4: Data Protection
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "2. Data Protection", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "You're handling video files, possibly PII, and internal corporate content. Encryption is non-negotiable.",
             font_size=14, color=MID_GRAY)

# In-Transit box
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.8), TEAL_LIGHT, border_color=TEAL, border_width=Pt(1.5))
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5), Inches(0.4),
             "\U0001F512 Encryption In Transit", font_size=20, color=TEAL, bold=True)

transit_items = [
    "\u2022 TLS 1.2 / TLS 1.3 minimum on all connections",
    "\u2022 HTTPS enforced everywhere \u2014 reject non-TLS",
    "\u2022 Kaltura download URLs: signed + HTTPS",
    "\u2022 S3 transfers: AWS SDK handles TLS",
    "\u2022 Zoom upload: HTTPS multipart POST",
    "\u2022 Certificate validation enabled",
    "",
    "\u2022 NEVER download to local disk",
    "\u2022 Stream directly: Kaltura \u2192 S3 \u2192 Zoom",
]
add_bullet_list(slide, Inches(1.2), Inches(2.8), Inches(5), Inches(3.5), transit_items, font_size=13)

# At-Rest box
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.8), LIGHT_GRAY, border_color=RGBColor(0xE5, 0xE7, 0xEB), border_width=Pt(1.5))
add_text_box(slide, Inches(7.4), Inches(2.2), Inches(5), Inches(0.4),
             "\U0001F5C4 Encryption At Rest", font_size=20, color=DARK, bold=True)

rest_items = [
    "S3 Bucket:",
    "\u2022 Server-Side Encryption (SSE-KMS)",
    "\u2022 Customer-managed KMS key",
    "\u2022 Key rotation annually",
    "\u2022 Block ALL public access",
    "\u2022 Block public ACLs",
    "",
    "DynamoDB:",
    "\u2022 Encryption at rest (KMS-managed)",
    "",
    "Secrets Manager:",
    "\u2022 Auto-encrypted via KMS",
]
add_bullet_list(slide, Inches(7.4), Inches(2.8), Inches(5), Inches(3.5), rest_items, font_size=13)


# ═══════════════════════════════════════════════════════════
# SLIDE 5: Network Security & Access Control
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "3. Network Security & Access Control", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)

# Network column
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.4),
             "\U0001F310 Network Security", font_size=20, color=TEAL, bold=True)

net_basic = add_shape(slide, Inches(0.8), Inches(2.1), Inches(5.8), Inches(2.2), LIGHT_GRAY)
add_text_box(slide, Inches(1.1), Inches(2.2), Inches(5.2), Inches(0.3),
             "Basic (Current)", font_size=14, color=DARK, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(2.6), Inches(5.2), Inches(1.5), [
    "\u2022 No public S3 buckets",
    "\u2022 Block public ACLs",
    "\u2022 Signed URLs for temporary access",
    "\u2022 HTTPS-only API endpoints",
], font_size=12)

net_ent = add_shape(slide, Inches(0.8), Inches(4.5), Inches(5.8), Inches(2.5), TEAL_LIGHT, border_color=TEAL, border_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(4.6), Inches(5.2), Inches(0.3),
             "Enterprise (Target)", font_size=14, color=TEAL, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(5.0), Inches(5.2), Inches(1.8), [
    "\u2022 Deploy Lambdas inside VPC",
    "\u2022 VPC endpoints for S3 + DynamoDB",
    "\u2022 No public internet exposure",
    "\u2022 AWS WAF in front of dashboard",
    "\u2022 Private S3 endpoints only",
], font_size=12)

# Access Control column
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4),
             "\U0001F6E1 Dashboard Access Control", font_size=20, color=DARK, bold=True)

ac_min = add_shape(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(2.2), LIGHT_GRAY)
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5.2), Inches(0.3),
             "Minimum Required", font_size=14, color=DARK, bold=True)
add_bullet_list(slide, Inches(7.3), Inches(2.6), Inches(5.2), Inches(1.5), [
    "\u2022 JWT-based authentication",
    "\u2022 Admin-only login required",
    "\u2022 Role-based access (Admin/Operator/Viewer)",
    "\u2022 Audit logs of who triggered migrations",
], font_size=12)

ac_ent = add_shape(slide, Inches(7.0), Inches(4.5), Inches(5.8), Inches(2.5), TEAL_LIGHT, border_color=TEAL, border_width=Pt(1))
add_text_box(slide, Inches(7.3), Inches(4.6), Inches(5.2), Inches(0.3),
             "Enterprise Grade", font_size=14, color=TEAL, bold=True)
add_bullet_list(slide, Inches(7.3), Inches(5.0), Inches(5.2), Inches(1.8), [
    "\u2022 SSO via SAML 2.0",
    "\u2022 SCIM provisioning",
    "\u2022 Enforced MFA for all users",
    "\u2022 Auth0 / Clerk / Supabase Auth",
    "\u2022 Session timeout policies",
], font_size=12)


# ═══════════════════════════════════════════════════════════
# SLIDE 6: Rate Limiting & Blast Radius
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "4. Rate Limiting & Blast Radius Control", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)

# Rate limiting
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.4),
             "\u26A1 Rate Limiting & Abuse Protection", font_size=20, color=TEAL, bold=True)

add_shape(slide, Inches(0.8), Inches(2.1), Inches(5.8), Inches(4.8), LIGHT_GRAY)
rl_items = [
    "Zoom Rate Limits:",
    "\u2022 Light: 30 req/sec (Pro), 80 (Business+)",
    "\u2022 Medium: 20 req/sec",
    "\u2022 Heavy: 10 req/sec",
    "",
    "Protocols:",
    "\u2022 Exponential backoff: 1s \u2192 2s \u2192 4s \u2192 8s \u2192 fail",
    "\u2022 Respect Retry-After header",
    "\u2022 Max concurrency cap (100 workers)",
    "\u2022 Circuit breaker on repeated 429s",
    "",
    "Kaltura:",
    "\u2022 Rate limits undocumented",
    "\u2022 Paginate conservatively",
    "\u2022 Add request throttling",
]
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.2), Inches(4.3), rl_items, font_size=13)

# Blast radius
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.4),
             "\U0001F6E1 Isolation & Blast Radius", font_size=20, color=DARK, bold=True)

add_shape(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.8), LIGHT_GRAY)
br_items = [
    "One bad video must NOT:",
    "\u2022 Crash the whole batch",
    "\u2022 Leak credentials",
    "\u2022 Corrupt migration state",
    "",
    "Best Practices:",
    "\u2022 Each video processed independently",
    "\u2022 Idempotent uploads (safe retries)",
    "\u2022 DynamoDB checkpointing per video",
    "\u2022 Dead-letter queue for failed jobs",
    "",
    "AWS Implementation:",
    "\u2022 SQS + DLQ for failed videos",
    "\u2022 Step Functions with Catch blocks",
    "\u2022 Independent error handling per video",
]
add_bullet_list(slide, Inches(7.3), Inches(2.3), Inches(5.2), Inches(4.3), br_items, font_size=13)


# ═══════════════════════════════════════════════════════════
# SLIDE 7: Audit & Logging
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "5. Audit, Logging & Incident Response", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
             "For enterprise customers, structured audit logs are a hard requirement.", font_size=14, color=MID_GRAY)

# Required log fields table
add_shape(slide, Inches(0.8), Inches(2.0), Inches(6.5), Inches(0.5), TEAL)
add_text_box(slide, Inches(1.0), Inches(2.05), Inches(6), Inches(0.4),
             "Every migration event must log:", font_size=14, color=WHITE, bold=True)

log_fields = [
    ("Video ID", "Unique source identifier"),
    ("Operator ID", "Who triggered the migration"),
    ("Source Timestamp", "When extracted from Kaltura"),
    ("Destination Timestamp", "When uploaded to Zoom"),
    ("Status", "pending / completed / failed"),
    ("Error Message", "Detailed error if failed"),
    ("SHA256 Checksum", "File integrity verification"),
    ("API Response Code", "HTTP status from each service"),
    ("File Size", "Bytes transferred"),
]

for i, (field, desc) in enumerate(log_fields):
    y = Inches(2.6) + i * Inches(0.42)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_shape(slide, Inches(0.8), y, Inches(6.5), Inches(0.4), bg_color)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.3),
                 field, font_size=12, color=DARK, bold=True)
    add_text_box(slide, Inches(3.5), y + Inches(0.05), Inches(3.5), Inches(0.3),
                 desc, font_size=12, color=MID_GRAY)

# Incident Response
add_text_box(slide, Inches(7.8), Inches(2.0), Inches(5), Inches(0.4),
             "\U0001F6A8 Incident Response Protocol", font_size=18, color=RED, bold=True)

add_shape(slide, Inches(7.8), Inches(2.6), Inches(5), Inches(4.2), RGBColor(0xFE, 0xF2, 0xF2), border_color=RED, border_width=Pt(1))
ir_items = [
    "Alerts & Monitoring:",
    "\u2022 CloudWatch alarms on high error rate",
    "\u2022 SNS notification on unauthorized access",
    "\u2022 Alert on IAM policy changes",
    "\u2022 Failed upload threshold alerts",
    "",
    "Response Procedures:",
    "\u2022 Credential rotation playbook",
    "\u2022 Data deletion procedure",
    "\u2022 Access revocation process",
    "\u2022 Post-incident review template",
    "",
    "Infrastructure:",
    "\u2022 AWS CloudTrail enabled",
    "\u2022 Immutable logging",
    "\u2022 90+ day log retention",
]
add_bullet_list(slide, Inches(8.0), Inches(2.8), Inches(4.6), Inches(3.8), ir_items, font_size=12)


# ═══════════════════════════════════════════════════════════
# SLIDE 8: Data Handling Protocol
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "6. Data Handling & Transfer Protocol", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)

# Flow diagram as cards
steps = [
    ("1", "Retrieve\nMetadata", "Kaltura API\nmedia.get()", RGBColor(0x7C, 0x3A, 0xED)),
    ("2", "Stream\nBinary", "Direct to S3\nNo local disk", RGBColor(0xD9, 0x77, 0x06)),
    ("3", "Verify\nChecksum", "SHA256 hash\nSize validation", TEAL),
    ("4", "Upload\nto Zoom", "HTTPS POST\nMultipart", BLUE),
    ("5", "Delete\nTemp File", "S3 lifecycle\nAuto-cleanup", GREEN),
    ("6", "Log\nResult", "DynamoDB +\nCloudWatch", DARK),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5) + i * Inches(2.1)
    y = Inches(1.6)

    # Card
    add_shape(slide, x, y, Inches(1.9), Inches(2.2), LIGHT_GRAY, border_color=color, border_width=Pt(2))

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), y + Inches(0.15), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.1), y + Inches(0.8), Inches(1.7), Inches(0.6),
                 title, font_size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.45), Inches(1.7), Inches(0.7),
                 desc, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(1.85), y + Inches(0.85), Inches(0.3), Inches(0.4),
                     "\u2192", font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Critical rules
add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8), RGBColor(0xFE, 0xF2, 0xF2), border_color=RED, border_width=Pt(1.5))
add_text_box(slide, Inches(1.1), Inches(4.35), Inches(11), Inches(0.4),
             "\u26A0  Critical Data Handling Rules", font_size=18, color=RED, bold=True)

rules = [
    "\u2022 NEVER store video files on local disk \u2014 stream directly between services",
    "\u2022 SHA256 checksum validation before AND after every transfer",
    "\u2022 S3 lifecycle rules: auto-delete staged files after successful migration",
    "\u2022 Verify file size matches Content-Length header",
    "\u2022 Private S3 bucket with NO public access, block public ACLs",
    "\u2022 Use presigned URLs with short expiry for any temporary access",
]
add_bullet_list(slide, Inches(1.1), Inches(4.9), Inches(11), Inches(2.0), rules, font_size=13)


# ═══════════════════════════════════════════════════════════
# SLIDE 9: Compliance
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "7. Compliance & Enterprise Readiness", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)

# Three compliance frameworks
frameworks = [
    ("SOC 2", "Financial Services", [
        "\u2022 Access control documentation",
        "\u2022 Change management process",
        "\u2022 Encryption at rest & in transit",
        "\u2022 Quarterly access reviews",
        "\u2022 Vendor risk assessment",
        "\u2022 Incident response plan",
    ], TEAL),
    ("GDPR", "EU Data Protection", [
        "\u2022 Data retention limits defined",
        "\u2022 Right to erasure supported",
        "\u2022 Data processing agreements",
        "\u2022 Cross-border transfer controls",
        "\u2022 Privacy impact assessment",
        "\u2022 Consent management",
    ], BLUE),
    ("HIPAA", "Healthcare (if applicable)", [
        "\u2022 PHI encryption mandatory",
        "\u2022 Access audit trails",
        "\u2022 Business associate agreements",
        "\u2022 Minimum necessary standard",
        "\u2022 Breach notification procedures",
        "\u2022 Security risk analysis",
    ], RGBColor(0x7C, 0x3A, 0xED)),
]

for i, (name, subtitle, items, color) in enumerate(frameworks):
    x = Inches(0.8) + i * Inches(4.2)
    y = Inches(1.6)

    # Card
    add_shape(slide, x, y, Inches(3.9), Inches(5.2), LIGHT_GRAY, border_color=color, border_width=Pt(2))

    # Header bar
    add_shape(slide, x, y, Inches(3.9), Inches(0.8), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.05), Inches(3.3), Inches(0.4),
                 name, font_size=24, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.45), Inches(3.3), Inches(0.3),
                 subtitle, font_size=12, color=RGBColor(0xD1, 0xFA, 0xE5))

    add_bullet_list(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.3), Inches(3.8), items, font_size=13)


# ═══════════════════════════════════════════════════════════
# SLIDE 10: Common Mistakes & Summary
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "8. Common Mistakes to Avoid", font_size=32, color=DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)

# Mistakes column
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(0xFE, 0xF2, 0xF2), border_color=RED, border_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5.2), Inches(0.4),
             "\u2718  Security Anti-Patterns", font_size=18, color=RED, bold=True)

mistakes = [
    "\u2022 Hardcoding API keys in Vercel env variables",
    "\u2022 Using AWS root credentials",
    "\u2022 Public S3 buckets",
    "\u2022 No token rotation",
    "\u2022 No retry / rate control",
    "\u2022 No audit logs",
    "\u2022 AdministratorAccess IAM policy",
    "\u2022 Storing videos on local disk",
    "\u2022 Open dashboard with no authentication",
    "\u2022 Using deprecated Zoom JWT apps",
]
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5.2), Inches(4.0), mistakes, font_size=14, color=RED)

# Checklist column
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(0xEC, 0xFD, 0xF5), border_color=GREEN, border_width=Pt(1))
add_text_box(slide, Inches(7.3), Inches(1.6), Inches(5.2), Inches(0.4),
             "\u2714  Minimum Secure MVP", font_size=18, color=GREEN, bold=True)

checklist = [
    "\u2714  OAuth 2.0 (Zoom S2S)",
    "\u2714  AWS Secrets Manager",
    "\u2714  IAM least privilege roles",
    "\u2714  Encrypted S3 (SSE-KMS)",
    "\u2714  Private bucket, no public access",
    "\u2714  JWT auth on dashboard",
    "\u2714  CloudWatch structured logging",
    "\u2714  Exponential backoff + retry",
    "\u2714  SHA256 file integrity checks",
    "\u2714  90-day log retention",
]
add_bullet_list(slide, Inches(7.3), Inches(2.2), Inches(5.2), Inches(4.0), checklist, font_size=14, color=GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE 11: Closing / CTA
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(6), Inches(0.6),
             "OpenExchange", font_size=20, color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(1.0),
             "From Demo to Enterprise Platform", font_size=40, color=WHITE, bold=True)

add_shape(slide, Inches(0.8), Inches(3.5), Inches(2), Inches(0.04), TEAL)

add_text_box(slide, Inches(0.8), Inches(4.0), Inches(10), Inches(1.2),
             "Without these security protocols, it's a demo tool.\n"
             "With them, it's a sellable migration platform.",
             font_size=22, color=RGBColor(0x94, 0xA3, 0xB8))

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.5),
             "Next Steps:  Security Audit  |  Minimal Secure MVP  |  Enterprise Hardening",
             font_size=16, color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(6), Inches(0.4),
             "openexc.com  |  Confidential", font_size=12, color=MID_GRAY)

# Decorative elements
add_shape(slide, Inches(10.5), Inches(1.5), Inches(2.5), Inches(5.5), RGBColor(0x00, 0x6B, 0x6E))
add_shape(slide, Inches(11.0), Inches(2.0), Inches(2.333), Inches(5.5), RGBColor(0x00, 0x55, 0x58))


# ── Save ──
output_path = "/Users/mithunmanjunatha/Desktop/claude 2/video-migration/OpenExchange_Security_Protocols.pptx"
prs.save(output_path)
print(f"\n  \u2705 Presentation saved to:\n  {output_path}\n")
print(f"  Slides: {len(prs.slides)}")
print(f"  Format: 16:9 widescreen")
print(f"  Branding: OpenExchange (Teal #008285)")
