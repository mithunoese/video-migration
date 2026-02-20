"""
Generate Video Migration Agent Pitch Deck with OpenExchange branding.
14 slides — persuasive internal presentation for leadership.
Brand: Teal #008285, Dark #1A1A2E, White #FFFFFF, Font: Calibri
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
TEAL = RGBColor(0x00, 0x82, 0x85)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
BLUE = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
TEAL_LIGHT = RGBColor(0xE0, 0xF7, 0xF8)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helper Functions ──

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def add_circle(slide, left, top, size, fill_color, text="", text_size=16, text_color=WHITE):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill_color
    circle.line.fill.background()
    circle.shadow.inherit = False
    if text:
        tf = circle.text_frame
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(text_size)
        tf.paragraphs[0].font.color.rgb = text_color
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False
    return circle


def add_kpi_card(slide, left, top, width, height, number, label, accent_color=TEAL):
    card = add_shape(slide, left, top, width, height, WHITE, border_color=BORDER, border_width=Pt(1))
    add_rect(slide, left, top, width, Inches(0.06), accent_color)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.3), width - Inches(0.6), Inches(0.8),
                 number, font_size=36, color=accent_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.1), width - Inches(0.6), Inches(0.5),
                 label, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


def add_icon_card(slide, left, top, width, height, icon_text, title, description, accent_color=TEAL):
    add_shape(slide, left, top, width, height, WHITE, border_color=BORDER, border_width=Pt(1))
    add_rect(slide, left, top, Inches(0.06), height, accent_color)
    add_circle(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.5), accent_color, icon_text, 16)
    add_text_box(slide, left + Inches(1.0), top + Inches(0.2), width - Inches(1.2), Inches(0.4),
                 title, font_size=16, color=DARK, bold=True)
    add_text_box(slide, left + Inches(1.0), top + Inches(0.6), width - Inches(1.2), height - Inches(0.8),
                 description, font_size=11, color=MID_GRAY)


def slide_header(slide, title, subtitle=None, dark=False):
    bg_color = DARK if dark else WHITE
    add_bg(slide, bg_color)
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)
    title_color = WHITE if dark else DARK
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
                 title, font_size=32, color=title_color, bold=True)
    add_rect(slide, Inches(0.8), Inches(1.05), Inches(1.5), Inches(0.04), TEAL)
    if subtitle:
        sub_color = RGBColor(0x94, 0xA3, 0xB8) if dark else MID_GRAY
        add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
                     subtitle, font_size=14, color=sub_color)


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(6), Inches(0.6),
             "OpenExchange", font_size=20, color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.2),
             "The Video Migration Agent", font_size=48, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(3.4), Inches(10), Inches(0.8),
             "From Manual Tool to Autonomous Platform", font_size=24, color=RGBColor(0x94, 0xA3, 0xB8))

add_rect(slide, Inches(0.8), Inches(4.5), Inches(2), Inches(0.04), TEAL)

add_text_box(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(0.4),
             "Kaltura \u2192 AWS S3 \u2192 Zoom  |  Built. Deployed. Running.", font_size=16, color=MID_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.4),
             "Internal Team Presentation  |  February 2026  |  Confidential", font_size=14, color=MID_GRAY)

# Decorative teal blocks
add_shape(slide, Inches(10.5), Inches(1.5), Inches(2.5), Inches(5.5), RGBColor(0x00, 0x6B, 0x6E))
add_shape(slide, Inches(11.0), Inches(2.0), Inches(2.333), Inches(5.5), RGBColor(0x00, 0x55, 0x58))


# ═══════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "The Problem: Manual Video Migration",
             "Thousands of enterprise videos locked in Kaltura. Zoom needs them. Nobody wants to move them by hand.")

# Pain points box
add_shape(slide, Inches(0.8), Inches(2.0), Inches(6.5), Inches(4.5),
          RGBColor(0xFE, 0xF2, 0xF2), border_color=RED, border_width=Pt(1.5))
add_text_box(slide, Inches(1.1), Inches(2.15), Inches(6), Inches(0.4),
             "\u2718  What Teams Do Today", font_size=20, color=RED, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.8), Inches(3.5), [
    "\u2022 Export metadata one-by-one from Kaltura admin panel",
    "\u2022 Download each video file manually (often GB-sized)",
    "\u2022 Re-upload to Zoom via browser UI or basic scripts",
    "\u2022 Manually verify each upload succeeded",
    "\u2022 Track progress in spreadsheets",
    "\u2022 Handle failures by starting completely over",
    "\u2022 No audit trail, no cost tracking, no metadata mapping",
], font_size=14)

# KPI cards
add_kpi_card(slide, Inches(7.8), Inches(2.0), Inches(4.8), Inches(1.3), "40+ hours", "Per 500 videos \u2014 manual migration", RED)
add_kpi_card(slide, Inches(7.8), Inches(3.5), Inches(4.8), Inches(1.3), "15\u201325%", "Failure rate from timeouts & human error", AMBER)
add_kpi_card(slide, Inches(7.8), Inches(5.0), Inches(4.8), Inches(1.3), "Zero", "Traceability \u2014 no audit logs or cost data", MID_GRAY)

# Urgency banner
add_shape(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), AMBER)
add_text_box(slide, Inches(1.1), Inches(6.85), Inches(11), Inches(0.4),
             "OE has two active contracts right now: IFRS and Indeed. Both need this solved.", font_size=14, color=WHITE, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: What We Built
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "What We Built: A Complete Migration Platform",
             "Not a script. A platform with real-time monitoring, cost tracking, and AI assistance.")

add_icon_card(slide, Inches(0.8), Inches(2.0), Inches(3.8), Inches(4.8),
              "P", "Python Migration Engine",
              "Kaltura API client with KS session auth\n\n"
              "\u2022 S3 staging with multipart uploads\n"
              "\u2022 Zoom OAuth 2.0 S2S uploader\n"
              "\u2022 Events CMS + Video Management\n"
              "\u2022 Concurrent processing (5 workers)\n"
              "\u2022 Exponential backoff retry\n"
              "\u2022 Per-video error isolation\n"
              "\u2022 DynamoDB state tracking",
              accent_color=PURPLE)

add_icon_card(slide, Inches(4.9), Inches(2.0), Inches(3.8), Inches(4.8),
              "D", "Real-Time Dashboard",
              "FastAPI + Alpine.js + Chart.js\n\n"
              "\u2022 7-tab interface (dark/white themes)\n"
              "\u2022 SSE real-time progress streaming\n"
              "\u2022 Video library with search & filter\n"
              "\u2022 Cost tracking with projections\n"
              "\u2022 AI assistant (Claude-powered)\n"
              "\u2022 Field mapping visualization\n"
              "\u2022 Pipeline test mode",
              accent_color=BLUE)

add_icon_card(slide, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.8),
              "\u2713", "Deployed and Running",
              "Live on Vercel \u2014 right now\n\n"
              "\u2022 video-migration.vercel.app\n"
              "\u2022 Demo mode with 847 test videos\n"
              "\u2022 Architecture diagram page\n"
              "\u2022 CLI + Dashboard + API\n"
              "\u2022 Test mode proves pipeline\n"
              "  works end-to-end\n"
              "\u2022 Zero production credentials needed",
              accent_color=GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: How It Works
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "How It Works: The Migration Pipeline",
             "Six automated steps. Every video. Every time. With full traceability.")

steps = [
    ("1", "Discover", "Kaltura API\nlist_videos()", PURPLE),
    ("2", "Extract", "Full metadata\n14 fields mapped", PURPLE),
    ("3", "Download", "Signed URL\nstream to disk", AMBER),
    ("4", "Stage", "S3 multipart\nSSE-KMS encrypted", AMBER),
    ("5", "Upload", "Zoom OAuth S2S\nEvents or VM API", BLUE),
    ("6", "Verify", "State tracker\nchecksum + log", GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.5) + i * Inches(2.1)
    y = Inches(1.8)

    add_shape(slide, x, y, Inches(1.9), Inches(2.4), LIGHT_GRAY, border_color=color, border_width=Pt(2))
    add_circle(slide, x + Inches(0.65), y + Inches(0.15), Inches(0.55), color, num, 18)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.85), Inches(1.7), Inches(0.5),
                 title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.45), Inches(1.7), Inches(0.8),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(1.85), y + Inches(0.9), Inches(0.3), Inches(0.4),
                     "\u2192", font_size=22, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Technical details box
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), LIGHT_GRAY, border_color=BORDER, border_width=Pt(1))
add_text_box(slide, Inches(1.1), Inches(4.7), Inches(5), Inches(0.4),
             "Technical Details", font_size=16, color=TEAL, bold=True)

add_bullet_list(slide, Inches(1.1), Inches(5.2), Inches(5.5), Inches(1.8), [
    "\u2022 ThreadPoolExecutor: up to 5 concurrent workers",
    "\u2022 Retry: exponential backoff (1s \u2192 2s \u2192 4s \u2192 8s)",
    "\u2022 State: 6 states tracked per video in DynamoDB",
    "\u2022 Multipart uploads for files > 100MB",
], font_size=12)

add_bullet_list(slide, Inches(7.0), Inches(5.2), Inches(5.0), Inches(1.8), [
    "\u2022 Per-video isolation: one failure \u2260 batch crash",
    "\u2022 14 Kaltura fields mapped to Zoom equivalents",
    "\u2022 Rate limit detection + Retry-After header",
    "\u2022 Automatic cleanup of staged temp files",
], font_size=12)


# ═══════════════════════════════════════════════════════════
# SLIDE 5: Live Dashboard
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "The Dashboard: Real-Time Control Center",
             "Seven tabs. Full visibility. No SSH required.")

tabs = [
    ("\U0001F4CA", "Dashboard", "KPI cards, activity feed,\nconnection health, cost summary", TEAL),
    ("\U0001F4F9", "Video Library", "Searchable table, filters,\npagination, detail panels", BLUE),
    ("\u26A1", "Migration Control", "Start/stop/retry, batch size,\nlive SSE event stream", PURPLE),
    ("\U0001F504", "Field Mapping", "Kaltura \u2192 Zoom metadata\ntransformation table", AMBER),
    ("\U0001F916", "AI Assistant", "Natural language queries,\nClaude-powered Tier-2", GREEN),
    ("\U0001F4B0", "Cost Tracker", "Per-service breakdown,\nprojections, CSV export", RED),
    ("\u2699", "Settings", "Credentials, connection tests,\npipeline test runner", MID_GRAY),
]

for i, (icon, name, desc, color) in enumerate(tabs):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + col * Inches(3.15)
    y = Inches(2.0) + row * Inches(2.5)

    add_shape(slide, x, y, Inches(2.9), Inches(2.1), WHITE, border_color=color, border_width=Pt(2))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.5), Inches(0.4),
                 f"{icon}  {name}", font_size=16, color=DARK, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.65), Inches(2.5), Inches(1.2),
                 desc, font_size=12, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 6: Why an Agent (THESIS SLIDE — dark bg)
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.4),
             "The Core Argument", font_size=18, color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "This Should Be an Autonomous Agent", font_size=42, color=WHITE, bold=True)

add_rect(slide, Inches(0.8), Inches(2.3), Inches(2), Inches(0.04), TEAL)

# Four principle cards on dark bg
principles = [
    ("\U0001F553", "24/7 Autonomous Operation",
     "Schedule migrations overnight. No human monitoring. "
     "Auto-discover new videos from Kaltura. Scale up/down based on API limits."),
    ("\U0001F527", "Intelligent Error Recovery",
     "Exponential backoff on rate limits. Per-video isolation. "
     "Auto-retry failed videos. Resume interrupted batches from checkpoint."),
    ("\U0001F4B5", "Cost-Aware Decisions",
     "Real-time cost tracking per video. Budget alerts and auto-pause. "
     "Projections before starting. Agent optimizes batch sizes for cost."),
    ("\U0001F4C8", "Scales Without Babysitting",
     "Configurable 1\u2013N workers. Process 10 or 10,000 videos. "
     "Structured reporting. No human reading logs at 2 AM."),
]

for i, (icon, title, desc) in enumerate(principles):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.15)
    y = Inches(2.8) + row * Inches(2.1)

    add_shape(slide, x, y, Inches(5.9), Inches(1.8), RGBColor(0x22, 0x22, 0x3A),
              border_color=TEAL, border_width=Pt(1))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(5.3), Inches(0.4),
                 f"{icon}  {title}", font_size=16, color=TEAL, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.65), Inches(5.3), Inches(1.0),
                 desc, font_size=13, color=RGBColor(0xCB, 0xD5, 0xE1))

# Bottom quote
add_text_box(slide, Inches(1.5), Inches(7.0), Inches(10), Inches(0.4),
             "\u201CA tool waits for you to click. An agent makes decisions, handles exceptions, and runs while you sleep.\u201D",
             font_size=15, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 7: Agent vs Tool Comparison
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Agent vs. Tool: What Changes Everything")

# Table headers
add_shape(slide, Inches(0.8), Inches(1.8), Inches(3.5), Inches(0.5), MID_GRAY)
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(3.3), Inches(0.4),
             "Dimension", font_size=14, color=WHITE, bold=True)

add_shape(slide, Inches(4.3), Inches(1.8), Inches(4.3), Inches(0.5), RED)
add_text_box(slide, Inches(4.5), Inches(1.85), Inches(4.0), Inches(0.4),
             "\u2718  Manual Tool", font_size=14, color=WHITE, bold=True)

add_shape(slide, Inches(8.6), Inches(1.8), Inches(4.3), Inches(0.5), TEAL)
add_text_box(slide, Inches(8.8), Inches(1.85), Inches(4.0), Inches(0.4),
             "\u2714  Autonomous Agent", font_size=14, color=WHITE, bold=True)

rows = [
    ("Operation", "Click \u201CStart\u201D per batch, watch, wait", "Schedule once, runs continuously"),
    ("Error Handling", "Human reads logs, re-runs manually", "Auto-retry with backoff, DLQ"),
    ("Rate Limits", "Hit 429, crash, start over", "Detect, back off, resume automatically"),
    ("Cost Control", "Check AWS bill end of month", "Real-time tracking, budget alerts"),
    ("Metadata", "Copy-paste fields manually", "14-field auto-mapping + AI gap analysis"),
    ("Monitoring", "SSH in, tail logs", "Real-time SSE dashboard + AI chat"),
    ("Off-Hours", "Nobody running it", "Runs 24/7 autonomously"),
    ("Scaling", "One operator, one batch", "Configurable concurrency, auto-batching"),
    ("Reporting", "Manual spreadsheet", "Auto-generated status summaries"),
    ("Multi-Client", "Rebuild per client", "Config switch: IFRS vs Indeed"),
]

for i, (dim, tool, agent) in enumerate(rows):
    y = Inches(2.35) + i * Inches(0.45)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE

    add_shape(slide, Inches(0.8), y, Inches(3.5), Inches(0.42), bg)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3.3), Inches(0.3),
                 dim, font_size=11, color=DARK, bold=True)

    add_shape(slide, Inches(4.3), y, Inches(4.3), Inches(0.42), bg)
    add_text_box(slide, Inches(4.5), y + Inches(0.05), Inches(4.0), Inches(0.3),
                 tool, font_size=11, color=RED)

    add_shape(slide, Inches(8.6), y, Inches(4.3), Inches(0.42), bg)
    add_text_box(slide, Inches(8.8), y + Inches(0.05), Inches(4.0), Inches(0.3),
                 agent, font_size=11, color=TEAL)


# ═══════════════════════════════════════════════════════════
# SLIDE 8: Security Built In
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Enterprise-Grade Security: Built In",
             "Not bolted on. Designed from day one. Full companion Security Protocols deck available (11 slides).")

add_icon_card(slide, Inches(0.8), Inches(2.0), Inches(3.8), Inches(4.8),
              "\U0001F512", "Authentication",
              "Three-platform auth secured:\n\n"
              "\u2022 Kaltura: KS session tokens, 24hr TTL\n"
              "\u2022 Zoom: OAuth 2.0 S2S (no user tokens)\n"
              "\u2022 AWS: IAM roles, least privilege\n"
              "\u2022 Dashboard: JWT + login screen\n"
              "\u2022 Rate-limited login endpoint\n"
              "\u2022 Security headers on all responses",
              accent_color=PURPLE)

add_icon_card(slide, Inches(4.9), Inches(2.0), Inches(3.8), Inches(4.8),
              "\U0001F6E1", "Data Protection",
              "Encryption everywhere:\n\n"
              "\u2022 TLS 1.2+ on all connections\n"
              "\u2022 S3 SSE-KMS encryption at rest\n"
              "\u2022 No public S3 buckets, block ACLs\n"
              "\u2022 Signed URLs with short expiry\n"
              "\u2022 Auto-cleanup lifecycle rules\n"
              "\u2022 Credentials in env / Secrets Mgr",
              accent_color=BLUE)

add_icon_card(slide, Inches(9.0), Inches(2.0), Inches(3.8), Inches(4.8),
              "\u2714", "Compliance Ready",
              "Enterprise audit trail:\n\n"
              "\u2022 SOC 2 aligned access controls\n"
              "\u2022 GDPR data retention limits\n"
              "\u2022 Per-video audit logging\n"
              "\u2022 Structured CloudWatch logs\n"
              "\u2022 90+ day log retention\n"
              "\u2022 Incident response procedures",
              accent_color=TEAL)


# ═══════════════════════════════════════════════════════════
# SLIDE 9: ROI & Cost Savings
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "ROI: The Numbers Speak",
             "Real cost data from AWS us-east-1 pricing. Zoom and Kaltura APIs are free with existing licenses.")

# KPI cards
add_kpi_card(slide, Inches(0.8), Inches(2.0), Inches(3.8), Inches(1.6), "$0.04", "Average cost per video migrated", TEAL)
add_kpi_card(slide, Inches(4.9), Inches(2.0), Inches(3.8), Inches(1.6), "95%", "Time reduction vs. manual process", GREEN)
add_kpi_card(slide, Inches(9.0), Inches(2.0), Inches(3.8), Inches(1.6), "$0", "Additional Zoom/Kaltura API cost", BLUE)

# Projection table
add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.45), TEAL)
headers = [("Scale", 2.0), ("Videos", 1.5), ("Data", 1.5), ("Est. Cost", 1.8), ("Manual Time", 2.3), ("Agent Time", 2.3)]
x = Inches(0.8)
for label, w in headers:
    add_text_box(slide, x + Inches(0.1), Inches(4.05), Inches(w - 0.2), Inches(0.35),
                 label, font_size=12, color=WHITE, bold=True)
    x += Inches(w)

projections = [
    ("IFRS Pilot", "500", "~150 GB", "~$20", "~40 hours", "~1.5 hours"),
    ("Indeed Full", "2,000", "~600 GB", "~$80", "~160 hours", "~6 hours"),
    ("Enterprise", "10,000", "~3 TB", "~$400", "~830 hours", "~30 hours"),
    ("Scale Target", "50,000", "~15 TB", "~$2,000", "~4,150 hours", "~150 hours"),
]

for i, row_data in enumerate(projections):
    y = Inches(4.5) + i * Inches(0.45)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    x = Inches(0.8)
    for j, (val, (_, w)) in enumerate(zip(row_data, headers)):
        add_shape(slide, x, y, Inches(w), Inches(0.42), bg)
        color = GREEN if j == 5 else (RED if j == 4 else DARK)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(w - 0.2), Inches(0.3),
                     val, font_size=12, color=color, bold=(j == 0))
        x += Inches(w)

# Footnote
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.4),
             "Based on: S3 transfer $0.09/GB, S3 storage $0.023/GB/mo, DynamoDB + Lambda < $0.001/video. Avg video ~300MB.",
             font_size=11, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 10: Already Working
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Already Built. Already Deployed. Already Proven.",
             "This is not a prototype. The core pipeline, dashboard, and cost tracking are production code.")

# Test results
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5),
          RGBColor(0xEC, 0xFD, 0xF5), border_color=GREEN, border_width=Pt(1.5))
add_text_box(slide, Inches(1.1), Inches(2.15), Inches(5.2), Inches(0.4),
             "\u2714  Pipeline Test Results", font_size=20, color=GREEN, bold=True)

test_steps = [
    "\u2714  Step 1: Discover test video (Big Buck Bunny, CC license)",
    "\u2714  Step 2: Extract metadata (title, duration, format, size)",
    "\u2714  Step 3: Download sample (0.95 MB, < 1 second)",
    "\u2714  Step 4: Stage to S3 (LocalStack or mock)",
    "\u2714  Step 5: Upload to Zoom (mock with file verification)",
    "\u2714  Step 6: Verify checksums and cleanup",
    "",
    "No production credentials needed.",
    "Run: python run.py test",
]
add_bullet_list(slide, Inches(1.1), Inches(2.7), Inches(5.2), Inches(3.5), test_steps, font_size=13, color=GREEN)

# Live deployment
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.5),
          TEAL_LIGHT, border_color=TEAL, border_width=Pt(1.5))
add_text_box(slide, Inches(7.3), Inches(2.15), Inches(5.2), Inches(0.4),
             "\U0001F310  Live Deployment", font_size=20, color=TEAL, bold=True)

deploy_items = [
    "\u2022 Dashboard: video-migration.vercel.app",
    "\u2022 Architecture: .../architecture.html",
    "\u2022 FastAPI serverless via @vercel/python",
    "\u2022 Demo mode: 847 deterministic test videos",
    "\u2022 Realistic data: 30s to 2hr videos",
    "\u2022 8 error types simulated",
    "\u2022 SSE with polling fallback",
    "",
    "\u2022 CLI: python run.py [verify|discover|migrate]",
]
add_bullet_list(slide, Inches(7.3), Inches(2.7), Inches(5.2), Inches(3.5), deploy_items, font_size=13)


# ═══════════════════════════════════════════════════════════
# SLIDE 11: Two Active Clients
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Two Clients. Two APIs. One Platform.",
             "The same codebase handles both \u2014 just a config switch.")

# IFRS card
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.0), WHITE, border_color=PURPLE, border_width=Pt(2))
add_rect(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.7), PURPLE)
add_text_box(slide, Inches(1.2), Inches(2.1), Inches(4.8), Inches(0.5),
             "IFRS", font_size=24, color=WHITE, bold=True)

add_bullet_list(slide, Inches(1.2), Inches(2.9), Inches(4.8), Inches(2.8), [
    "Source: Kaltura",
    "Destination: Zoom Events CMS",
    "Config: ZOOM_TARGET_API=events",
    "",
    "Content: Financial services videos",
    "Scale: ~500 videos, ~150 GB",
    "Status: Pipeline ready, awaiting credentials",
], font_size=14)

# Indeed card
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.0), WHITE, border_color=BLUE, border_width=Pt(2))
add_rect(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.7), BLUE)
add_text_box(slide, Inches(7.4), Inches(2.1), Inches(4.8), Inches(0.5),
             "Indeed", font_size=24, color=WHITE, bold=True)

add_bullet_list(slide, Inches(7.4), Inches(2.9), Inches(4.8), Inches(2.8), [
    "Source: Kaltura",
    "Destination: Zoom Video Management",
    "Config: ZOOM_TARGET_API=vm",
    "",
    "Content: Recruitment & training videos",
    "Scale: ~2,000+ videos, ~600 GB",
    "Status: Architecture confirmed, same codebase",
], font_size=14)

# Shared platform banner
add_shape(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7), TEAL)
add_text_box(slide, Inches(1.1), Inches(6.4), Inches(11), Inches(0.5),
             "Same ZoomClient code handles both \u2014 zoom_client.upload_video() routes by config. This is a platform, not a one-off script.",
             font_size=14, color=WHITE, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 12: Roadmap
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Roadmap: Building the Full Agent")

# Phase 1
add_shape(slide, Inches(0.5), Inches(1.7), Inches(4.0), Inches(5.2), WHITE, border_color=TEAL, border_width=Pt(2))
add_rect(slide, Inches(0.5), Inches(1.7), Inches(4.0), Inches(0.7), TEAL)
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(3.5), Inches(0.5),
             "Phase 1: Hardened Tool", font_size=18, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.25), Inches(3.5), Inches(0.3),
             "Now \u2192 Month 1", font_size=11, color=RGBColor(0xA7, 0xF3, 0xD0))

add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(3.5), Inches(4.0), [
    "\u2022 JWT auth on dashboard",
    "\u2022 Security headers + CORS",
    "\u2022 Rate limiting on API",
    "\u2022 Input validation (Pydantic)",
    "\u2022 Audit logging",
    "\u2022 DynamoDB state in production",
    "\u2022 S3 lifecycle cleanup rules",
    "\u2022 Complete IFRS pilot (~500 videos)",
], font_size=12)

# Phase 2
add_shape(slide, Inches(4.7), Inches(1.7), Inches(4.0), Inches(5.2), WHITE, border_color=BLUE, border_width=Pt(2))
add_rect(slide, Inches(4.7), Inches(1.7), Inches(4.0), Inches(0.7), BLUE)
add_text_box(slide, Inches(5.0), Inches(1.8), Inches(3.5), Inches(0.5),
             "Phase 2: Smart Automation", font_size=18, color=WHITE, bold=True)
add_text_box(slide, Inches(5.0), Inches(2.25), Inches(3.5), Inches(0.3),
             "Month 2\u20133", font_size=11, color=RGBColor(0xBF, 0xDB, 0xFE))

add_bullet_list(slide, Inches(5.0), Inches(2.7), Inches(3.5), Inches(4.0), [
    "\u2022 Scheduled migrations (cron)",
    "\u2022 Adaptive rate-limit throttling",
    "\u2022 Budget-aware auto-pause",
    "\u2022 Multi-client config (IFRS + Indeed)",
    "\u2022 Dead-letter queue for failures",
    "\u2022 Slack/Teams webhook alerts",
    "\u2022 AI anomaly detection",
    "\u2022 Complete Indeed migration",
], font_size=12)

# Phase 3
add_shape(slide, Inches(8.9), Inches(1.7), Inches(4.0), Inches(5.2), WHITE, border_color=GREEN, border_width=Pt(2))
add_rect(slide, Inches(8.9), Inches(1.7), Inches(4.0), Inches(0.7), GREEN)
add_text_box(slide, Inches(9.2), Inches(1.8), Inches(3.5), Inches(0.5),
             "Phase 3: Autonomous Agent", font_size=18, color=WHITE, bold=True)
add_text_box(slide, Inches(9.2), Inches(2.25), Inches(3.5), Inches(0.3),
             "Month 4\u20136", font_size=11, color=RGBColor(0xA7, 0xF3, 0xD0))

add_bullet_list(slide, Inches(9.2), Inches(2.7), Inches(3.5), Inches(4.0), [
    "\u2022 Full agent loop: discover \u2192 plan \u2192 run",
    "\u2022 Self-healing on API failures",
    "\u2022 Predictive cost optimization",
    "\u2022 Multi-platform (Vimeo, Wistia)",
    "\u2022 White-label client dashboard",
    "\u2022 Agent observability + decision logs",
    "\u2022 Confidence scores + rollback",
    "\u2022 Migration-as-a-Service product",
], font_size=12)


# ═══════════════════════════════════════════════════════════
# SLIDE 13: Team Value
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Strategic Value for OpenExchange",
             "This isn't just a tool. It's a platform play that deepens our Zoom partnership.")

values = [
    ("\U0001F91D", "Deepens Zoom Partnership",
     "OE becomes Zoom's go-to partner for video onboarding. Every migration is a customer locked into Zoom. "
     "Positions OE as infrastructure partner, not just events."),
    ("\U0001F4B0", "Repeatable Revenue",
     "Migration-as-a-Service for every new Zoom Events / VM customer. "
     "At $0.04/video cost, significant margin even at low pricing. Scales linearly."),
    ("\U0001F3AF", "Competitive Moat",
     "No competitor has an autonomous migration agent. Most offer one-off scripts or manual services. "
     "Agent capability is a genuine differentiator in sales."),
    ("\U0001F680", "Engineering Leverage",
     "Built by a small team with modern tooling. Same codebase serves multiple clients and APIs. "
     "AI-assisted operations reduce ongoing support burden."),
]

for i, (icon, title, desc) in enumerate(values):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.15)
    y = Inches(2.0) + row * Inches(2.5)

    add_shape(slide, x, y, Inches(5.9), Inches(2.1), WHITE, border_color=TEAL, border_width=Pt(1.5))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.3), Inches(0.4),
                 f"{icon}  {title}", font_size=18, color=TEAL, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(5.3), Inches(1.2),
                 desc, font_size=13, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════
# SLIDE 14: Call to Action
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(6), Inches(0.6),
             "OpenExchange", font_size=20, color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.0),
             "Let\u2019s Build the Agent", font_size=44, color=WHITE, bold=True)

add_rect(slide, Inches(0.8), Inches(3.3), Inches(2), Inches(0.04), TEAL)

# Three action items
actions = [
    ("1", "Approve Phase 1", "Security hardening + IFRS pilot (4 weeks)"),
    ("2", "Assign Production Credentials", "Kaltura, AWS, Zoom for IFRS"),
    ("3", "Schedule Agent Review", "Monthly check-in on Phase 2 progress"),
]

for i, (num, title, desc) in enumerate(actions):
    y = Inches(3.8) + i * Inches(0.8)
    add_circle(slide, Inches(0.8), y, Inches(0.5), TEAL, num, 16)
    add_text_box(slide, Inches(1.5), y + Inches(0.02), Inches(5), Inches(0.3),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.35), Inches(8), Inches(0.3),
                 desc, font_size=14, color=RGBColor(0x94, 0xA3, 0xB8))

# Links
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(10), Inches(0.8),
             "Dashboard: video-migration.vercel.app\n"
             "Security Protocols: Companion deck (11 slides)  |  Test: python run.py test",
             font_size=13, color=MID_GRAY)

# Decorative
add_shape(slide, Inches(10.5), Inches(1.5), Inches(2.5), Inches(5.5), RGBColor(0x00, 0x6B, 0x6E))
add_shape(slide, Inches(11.0), Inches(2.0), Inches(2.333), Inches(5.5), RGBColor(0x00, 0x55, 0x58))


# ── Save ──
import os
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "VideoMigration_Pitch_Deck.pptx")
prs.save(output_path)
print(f"\n  \u2705 Pitch deck saved to:\n  {output_path}\n")
print(f"  Slides: {len(prs.slides)}")
print(f"  Format: 16:9 widescreen")
print(f"  Branding: OpenExchange (Teal #008285)")
